import numpy as np
import pandas as pd
import regex as re
import io
import fitz  #PyMuPDF for PDFs
import pytesseract
import cv2
from PIL import Image
from fastapi import UploadFile
from app.utils.external_api import call_gemini_api, analyze_image
from docx import Document
from pptx import Presentation
import logging

async def process_document(file: UploadFile, summary_length: int, focus_sections: str, language: str):
    try:
        filename = file.filename
        file_extension = filename.split(".")[-1].lower() if filename else ""
        text_content = ""
        image_text_content = ""
        image_count = 0
        image_analysis_results = []

        if file_extension == "pdf":
            text_content, image_text_content, image_count, image_analysis_results = await extract_text_from_pdf(file)
        elif file_extension == "docx":
            text_content = await extract_text_from_docx(file)
        elif file_extension == "pptx":
            text_content = await extract_text_from_pptx(file)
        elif file_extension == "xlsx":
            text_content = await extract_text_from_xlsx(file)
        else:
            return {"summary": "Unsupported file format", "images": 0}

        # Combine extracted text, OCR text, and image analysis results
        combined_text = (text_content + "\n" + image_text_content + "\n" + "\n".join(image_analysis_results)).strip()

        if not combined_text:
            return {"summary": "No content to summarize", "images": image_count}

        if focus_sections:
            combined_text = filter_focus_sections(combined_text, focus_sections)

        # Generate the final summary using Gemini API
        summary = await summarize_text(combined_text, summary_length, language)
        return {
            "summary": summary,
            "images": image_count
        }
    except Exception as e:
        return {"summary": f"Error processing document: {str(e)}", "images": 0}

async def extract_text_from_pdf(file):
    """Extracts text and images from a PDF file and analyzes images using Gemini API."""
    text_content = ""
    image_text_content = ""
    image_count = 0
    image_analysis_results = []

    file_bytes = await file.read()
    pdf_document = fitz.open(stream=file_bytes, filetype="pdf")

    for page in pdf_document:
        text_content += page.get_text("text") + "\n"

        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = pdf_document.extract_image(xref)
            image_bytes = base_image["image"]
            image = Image.open(io.BytesIO(image_bytes))

            # Preprocess image before OCR
            processed_image = preprocess_image(image)
            
            # Extract text from the image using OCR
            extracted_text = pytesseract.image_to_string(processed_image)
            extracted_text = clean_extracted_text(extracted_text)
            image_text_content += extracted_text + "\n"
            image_count += 1

            # Send image separately to Gemini API
            image_gemini_text = await analyze_image(image_bytes)
            image_analysis_results.append(image_gemini_text)

    return text_content.strip(), image_text_content.strip(), image_count, image_analysis_results

async def extract_text_from_docx(file):
    file_content = await file.read()
    doc = Document(io.BytesIO(file_content))  #Convert to BytesIO before passing to Document

    text = "\n".join([para.text for para in doc.paragraphs])
    
    return text

async def extract_text_from_pptx(file):
    """Extracts text from a PowerPoint (.pptx) file."""
    file_content = await file.read()
    ppt = Presentation(io.BytesIO(file_content))

    text_content = ""
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content += shape.text + "\n"
    
    return text_content.strip()

async def extract_text_from_xlsx(file):
    """Extracts text from an Excel (.xlsx) file."""
    file_content = await file.read()
    excel = pd.ExcelFile(io.BytesIO(file_content))

    text_content = ""
    for sheet_name in excel.sheet_names:
        df = excel.parse(sheet_name, dtype=str).fillna("")
        for row in df.itertuples(index=False, name=None):
            text_content += " ".join(row) + "\n"
    
    return text_content.strip()

#FUNCTION to summarize
async def summarize_text(text: str, summary_length: int, language: str) -> str:
    """
    Uses the Gemini API to analyze the document text and generate a summary.
    Handles long texts by processing in chunks.
    """
    chunk_size = 5000
    chunks = [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]
    summaries = [
        await call_gemini_api(f"Summarize the following document in {language}. Keep it approximately {summary_length} words:\n\n{chunk}")
        for chunk in chunks
    ]
    return "\n".join(summaries)

#HELPER functions -

#This filter function is over doing stuff
#Will have to come up with a more suitable function
def filter_focus_sections(text: str, focus_sections: str) -> str:
    """
    Filters the document text based on the provided focus sections.
    """
    sections = focus_sections.split(",")
    filtered_text = ""

    for section in sections:
        section = section.strip()
        pattern = rf"{section}.*?(?=\n[A-Z])"
        matches = re.findall(pattern, text, re.DOTALL)

        if matches:
            filtered_text += " ".join(matches) + "\n"

    return filtered_text if filtered_text else text

def preprocess_image(image):
    #Convert to grayscale
    gray = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2GRAY)

    #Rescale the image to improve OCR accuracy
    scale_factor = 2  # Adjust as needed
    height, width = gray.shape[:2]
    resized = cv2.resize(gray, (width * scale_factor, height * scale_factor), interpolation=cv2.INTER_LINEAR)

    #Apply adaptive thresholding
    thresh = cv2.adaptiveThreshold(resized, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2)

    #Optional: Denoise the image
    denoised = cv2.fastNlMeansDenoising(thresh, None, 30, 7, 21)

    return Image.fromarray(denoised)

def clean_extracted_text(text):
    #Remove non-ASCII characters
    cleaned_text = re.sub(r'[^\x00-\x7F]+', ' ', text)

    #Remove extra whitespace
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text).strip()    
    return cleaned_text
